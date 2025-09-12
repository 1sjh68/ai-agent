# utils/file_handler.py

import os
import json
import logging
import fitz  # PyMuPDF
from datetime import datetime
import docx
import pptx
import io
import csv
from config import Config
# --- [核心修改] 将可选的导入包围在try...except块中 ---
TESSERACT_AVAILABLE = False
try:
    from PIL import Image
    from pdf2image import convert_from_path
    import pytesseract
    # 如果您在Windows上安装了Tesseract，请取消下面这行的注释，并确保路径正确
    # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    TESSERACT_AVAILABLE = True
    logging.info("Tesseract-OCR and pdf2image found. OCR functionality is enabled.")
except ImportError:
    logging.warning("Tesseract-OCR or pdf2image not found. The application will run without OCR capabilities.")
    logging.warning("To enable processing of scanned/image-based PDFs, please install Tesseract-OCR and the necessary Python libraries.")
# --- [修改结束] ---


def _read_txt(file_path: str) -> str:
    """读取 .txt 文件内容。"""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def _read_docx(file_path: str) -> str:
    """读取 .docx 文件内容。"""
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text])

def _read_pptx(file_path: str) -> str:
    """读取 .pptx 文件中所有幻灯片的文本内容。"""
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def _read_pdf_with_ocr_fallback(file_path: str) -> str:
    """
    (V2 - Tesseract可选版) 智能读取 .pdf 文件。
    首先尝试直接提取文本，如果文本量过少且Tesseract可用，则自动启用OCR。
    """
    text = ""
    try:
        page_count = 0
        with fitz.open(file_path) as doc:
            text = "".join(page.get_text() for page in doc).strip()
            page_count = len(doc)
        
        # 检查是否需要OCR
        is_scanned_pdf = len(text) < 50 * page_count

        if is_scanned_pdf and TESSERACT_AVAILABLE:
            logging.info(f"  - Direct text extraction yielded little text ({len(text)} chars). Attempting OCR...")
            try:
                images = convert_from_path(file_path)
                ocr_texts = [pytesseract.image_to_string(image, lang='chi_sim+eng') for image in images] # 尝试中英双语
                text = "\n".join(ocr_texts)
                logging.info(f"  - OCR successful, recognized {len(text)} characters.")
            except Exception as ocr_error:
                # [核心修改] 如果OCR过程失败（例如，没找到tesseract.exe），则记录错误并回退
                logging.error(f"  - OCR process failed: {ocr_error}")
                logging.error("  - This could be due to a missing Tesseract installation or incorrect path.")
                text += "\n\n[OCR FAILED: Could not process images in this PDF.]"
        elif is_scanned_pdf and not TESSERACT_AVAILABLE:
            # 如果需要OCR但Tesseract不可用
            logging.warning("  - This appears to be a scanned PDF, but Tesseract-OCR is not available. Skipping OCR.")
            text += "\n\n[OCR SKIPPED: Tesseract-OCR not found. Text from images in this PDF was not extracted.]"
        else:
            logging.info("  - PDF appears to be text-based. Direct extraction successful.")

    except Exception as e:
        logging.error(f"An error occurred while processing PDF file {file_path}: {e}", exc_info=True)
        return f"[Error: Could not read PDF file {os.path.basename(file_path)}]"
            
    return text

def load_external_data(config: Config, file_paths: list[str]) -> str:
    """
    从给定的文件路径列表加载所有文本内容。
    支持 .txt, .pdf (包括图片型), .docx, .pptx。
    """
    if not file_paths:
        return ""
    
    readers = {
        '.txt': _read_txt,
        '.pdf': _read_pdf_with_ocr_fallback,
        '.docx': _read_docx,
        '.pptx': _read_pptx,
        '.ppt': _read_pptx
    }
    
    all_content = []
    for fp in file_paths:
        if not fp or not os.path.exists(fp):
            logging.warning(f"External data file not found, skipped: {fp}")
            continue
        
        ext = os.path.splitext(fp)[1].lower()
        content = ""
        
        if ext in readers:
            try:
                logging.info(f"Reading {ext.upper()} file: {fp} ...")
                content = readers[ext](fp)
                logging.info(f"  - Successfully read {os.path.basename(fp)} ({len(content)} chars)")
                
                file_header = f"\n\n--- Start of file: {os.path.basename(fp)} ---\n"
                file_footer = f"\n--- End of file: {os.path.basename(fp)} ---\n\n"
                all_content.append(file_header + content + file_footer)
                
            except Exception as e:
                logging.error(f"Error reading external data file {fp}: {e}")
        else:
            logging.warning(f"Unsupported file type: {ext} (file: {fp}). Skipped.")
            continue
            
    return "\n".join(all_content)

def save_checkpoint(config: Config, iteration: int, solution: str, feedback_history: list,
                    initial_problem: str, initial_solution_target_chars: int,
                    max_iterations: int, external_data_checksum: str,
                    document_outline_data: dict | None = None,
                    successful_patches: list[dict] | None = None,
                    research_briefs_history: list[str] | None = None,
                    style_guide: str | None = None
                    ):
    checkpoint_data = {
        "metadata": { "version": "1.1", "timestamp": datetime.now().isoformat() },
        "state": {
            "iteration": iteration, "current_solution": solution,
            "feedback_history": feedback_history, "initial_problem": initial_problem,
            "initial_solution_target_chars": initial_solution_target_chars,
            "max_iterations": max_iterations, "external_data_checksum": external_data_checksum,
            "document_outline_data": document_outline_data,
            "successful_patches": successful_patches if successful_patches else [],
            "research_briefs_history": research_briefs_history if research_briefs_history else [],
            "style_guide": style_guide
        }
    }
    if not config.session_dir:
        logging.error("Session directory not set, cannot save checkpoint.")
        return
    path = os.path.join(config.session_dir, config.checkpoint_file_name)
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(checkpoint_data, f, ensure_ascii=False, indent=4)
        logging.info(f"\n--- Checkpoint saved to {path} (after iteration {iteration + 1}) ---")
    except Exception as e:
        logging.error(f"\nError saving checkpoint: {e}")

def load_checkpoint(config: Config):
    if not config.session_dir:
        logging.error("Session directory not set, cannot load checkpoint.")
        return None
    path = os.path.join(config.session_dir, config.checkpoint_file_name)
    if not os.path.exists(path):
        return None
    try:
        with open(path, "r", encoding="utf-8") as f:
            checkpoint_data = json.load(f)
        state = checkpoint_data.get("state", checkpoint_data)
        logging.info(f"\n--- Successfully loaded checkpoint from {path} (saved after iteration {state.get('iteration', -1) + 1}) ---")
        state.setdefault("successful_patches", [])
        state.setdefault("research_briefs_history", [])
        state.setdefault("style_guide", "")
        return state
    except (json.JSONDecodeError, KeyError) as e:
        logging.error(f"\nError loading checkpoint: {e}. File may be corrupt or incompatible, will be deleted.")
        try:
            os.remove(path)
        except Exception as remove_err:
            logging.error(f"Error deleting corrupt checkpoint file: {remove_err}")
    except Exception as e:
        logging.error(f"\nUnknown error loading checkpoint: {e}", exc_info=True)
    return None

def delete_checkpoint(config: Config):
    if not config.session_dir:
        return
    path = os.path.join(config.session_dir, config.checkpoint_file_name)
    if os.path.exists(path):
        try:
            os.remove(path)
            logging.info(f"\n--- Task complete, checkpoint file {path} deleted ---")
        except Exception as e:
            logging.error(f"\nError deleting checkpoint file: {e}")


def parse_and_validate_paths(path_string: str) -> list[str]:
    """
    (Moved from main.py)
    Parses a comma-separated string of file paths, validates their existence,
    and returns a list of valid paths. Handles various quoting and spacing issues.
    """
    if not path_string or not path_string.strip():
        return []
    # 纠正Windows路径并去除首尾空白和引号
    corrected_path_string = path_string.replace('\\', '/').strip()
    if corrected_path_string.startswith('"') and corrected_path_string.endswith('"'):
        corrected_path_string = corrected_path_string[1:-1]
    
    potential_paths = []
    try:
        # 使用csv模块来正确处理带引号和逗号的复杂路径
        string_reader = io.StringIO(corrected_path_string)
        path_reader = csv.reader(string_reader, delimiter=',', quotechar='"', skipinitialspace=True)
        potential_paths = next(path_reader)
    except (StopIteration, csv.Error) as e:
        logging.warning(f"CSV parsing of paths failed, falling back to simple split. Error: {e}")
        # 回退机制：对于简单的、不带引号的路径字符串
        potential_paths = [p.strip() for p in corrected_path_string.split(',') if p.strip()]

    valid_paths = []
    for path in potential_paths:
        # 清理每个独立路径的首尾空白和可能存在的引号
        clean_path = path.strip().strip('"')
        if os.path.exists(clean_path):
            valid_paths.append(clean_path)
        else:
            logging.warning(f"Provided external file path does not exist and was skipped: '{clean_path}'")
    return valid_paths